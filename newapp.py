import streamlit as st
import pandas as pd
import numpy as np
import os
import matplotlib.pyplot as plt
import google.generativeai as genai
from openai import OpenAI
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# ==========================================
# 1. 基础配置与全局变量
# ==========================================
st.set_page_config(page_title="AI 智能科研与商业数据穿透 Agent", layout="wide")

# 初始化所有的持久化状态
if "analyzed" not in st.session_state: st.session_state.analyzed = False
if "file_type" not in st.session_state: st.session_state.file_type = "" # 'text' 或 'dataset'
if "raw_text" not in st.session_state: st.session_state.raw_text = ""
if "outline" not in st.session_state: st.session_state.outline = ""
if "df_cleaned" not in st.session_state: st.session_state.df_cleaned = None
if "df_hierarchical" not in st.session_state: st.session_state.df_hierarchical = None
if "diagnostic_report" not in st.session_state: st.session_state.diagnostic_report = {}
if "chat_history_dataset" not in st.session_state: st.session_state.chat_history_dataset = []
if "chat_history_text" not in st.session_state: st.session_state.chat_history_text = []

# 自适应列名映射词条字典
COLUMN_MAP = {
    'product_id': ['Product_ID', 'StockCode', 'product_id', '产品ID', '商品编码', 'sku'],
    'category': ['Category', 'product_category_name', '产品品类', '品类', 'category_name_1'],
    'price': ['Final_Price(Rs.)', 'UnitPrice', 'Price (Rs.)', 'price', '单价', '销售额'],
    'sales': ['order_item_id', 'Quantity', '销量', '数量', 'qty_ordered']
}

# 设置图表中文字体
plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# ==========================================
# 2. 核心后端引擎：本地非 AI 纯计算模块
# ==========================================
def extract_text_from_doc(uploaded_file):
    """文本类文件纯本地解析"""
    file_ext = uploaded_file.name.split(".")[-1].lower()
    text = ""
    if file_ext == "pdf":
        reader = PdfReader(uploaded_file)
        for page in reader.pages:
            t = page.extract_text()
            if t: text += t + "\n"
    elif file_ext in ["docx", "doc"]:
        doc = Document(uploaded_file)
        for para in doc.paragraphs: text += para.text + "\n"
    elif file_ext in ["pptx", "ppt"]:
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text: text += shape.text + "\n"
    elif file_ext == "txt":
        text = uploaded_file.read().decode("utf-8", errors="ignore")
    return text

def process_dataset_locally(uploaded_file):
    """CSV 数据集本地标准化清洗与穿透分析逻辑 (不消耗 Key)"""
    df_raw = pd.read_csv(uploaded_file, low_memory=False)
    df_mapped = pd.DataFrame()
    
    # 智能映射
    for logic, candidates in COLUMN_MAP.items():
        match = [c for c in df_raw.columns if c in candidates]
        if match: df_mapped[logic] = df_raw[match[0]]
        
    # 基础清洗保护
    if 'category' not in df_mapped.columns: df_mapped['category'] = '未分类未知品类'
    else: df_mapped['category'] = df_mapped['category'].fillna('未分类未知品类')
    
    if 'price' not in df_mapped.columns: df_mapped['price'] = 0.0
    if df_mapped['price'].dtype == 'object':
        df_mapped['price'] = df_mapped['price'].astype(str).str.replace(r'[^\d.]', '', regex=True)
        df_mapped['price'] = pd.to_numeric(df_mapped['price'], errors='coerce')
    df_mapped['price'] = df_mapped['price'].fillna(0.0)
    
    if 'sales' in df_mapped.columns: df_mapped['actual_sales'] = df_mapped['sales']
    else: df_mapped['actual_sales'] = 1

    # 聚合大类与单品
    df_product = df_mapped.groupby(['category', 'product_id']).agg(
        total_revenue=('price', 'sum'), sales_count=('actual_sales', 'count' if 'sales' not in df_mapped.columns else 'sum')
    ).reset_index()
    
    df_category = df_product.groupby('category').agg(
        total_revenue=('total_revenue', 'sum'), sales_count=('sales_count', 'sum')
    ).reset_index().sort_values(by='total_revenue', ascending=False)
    
    # ABC 矩阵分级
    df_category['cum_pct'] = df_category['total_revenue'].cumsum() / df_category['total_revenue'].sum()
    df_category['rank'] = df_category['cum_pct'].apply(lambda x: 'A' if x <= 0.8 else ('B' if x <= 0.95 else 'C'))
    
    # 构建另起一行的多层级穿透大表
    hierarchical_rows = []
    for _, cat_row in df_category.iterrows():
        cat_name = cat_row['category']
        hierarchical_rows.append({
            '数据层级': '【品类大类】', '名称/编码ID': cat_name,
            '总销售额(利润)': cat_row['total_revenue'], '总销量': cat_row['sales_count'], '决策分级': cat_row['rank']
        })
        cat_products = df_product[df_product['category'] == cat_name].sort_values(by='total_revenue', ascending=False)
        for _, prod_row in cat_products.iterrows():
            hierarchical_rows.append({
                '数据层级': '  └─ 具体单品 SKU', '名称/编码ID': prod_row['product_id'],
                '总销售额(利润)': prod_row['total_revenue'], '总销量': prod_row['sales_count'], '决策分级': '单品穿透'
            })
            
    df_hierarchical = pd.DataFrame(hierarchical_rows)
    
    # 动态风险警示线逻辑
    a_cats = df_category[df_category['rank'] == 'A']
    c_cats = df_category[df_category['rank'] == 'C']
    avg_rev_A = a_cats['total_revenue'].mean() if not a_cats.empty else 1
    avg_rev_C = c_cats['total_revenue'].mean() if not c_cats.empty else 0
    ratio = avg_rev_C / avg_rev_A if avg_rev_A > 0 else 0
    
    if ratio <= 0.10:
        warning = f"【裁剪预警】C类均值仅为A类的 {ratio*100:.1f}%，大类分化严重，系统建议：您可以果断砍掉整个C类品类，将资源集中于核心业务。"
    elif ratio >= 0.90:
        warning = f"【结构预警】C类均值达A类的 {ratio*100:.1f}%，品类大类间大体平衡，不宜盲目砍掉整个大类，请转而执行底层低效单品的精准淘汰。"
    else:
        warning = f"【混合预警】C类均值为A类的 {ratio*100:.1f}%，处于中游，建议保持现有大类关注，并精准下架拖后腿的低效具体单品。"
        
    # 高阶统计学假设检验建议引擎 (满足要求4)
    stat_advise = "【统计学模型方向建议】：\n"
    numeric_cols = df_raw.select_dtypes(include=[np.number]).columns.tolist()
    
    if len(numeric_cols) >= 2:
        stat_advise += f"1. 检测到多个连续数字变量 {numeric_cols[:3]}，满足多元线性回归条件，建议构建 OLS 模型的方差矩阵。\n"
    if 'category' in df_mapped.columns and len(df_mapped['category'].unique()) == 2:
        stat_advise += "2. 当前分类标签特征呈二分类分布，建议对双特征组执行【独立样本 T 检验】来测度显著性差异。\n"
    elif 'category' in df_mapped.columns and len(df_mapped['category'].unique()) > 2:
        stat_advise += "3. 检测到分类标签多于两组，样本群体宏观相近，建议执行【单因素方差分析 (ANOVA)】来检验类别间效应。\n"

    report_dict = {
        "warning": warning, "stat_advise": stat_advise,
        "top_cat": df_category.head(3), "bottom_cat": df_category.tail(3),
        "top_prod": df_product.sort_values(by='total_revenue', ascending=False).head(3),
        "bottom_prod": df_product.sort_values(by='total_revenue', ascending=False).tail(3)
    }
    
    return df_category, df_hierarchical, report_dict

# ==========================================
# 3. AI 调用中转模块
# ==========================================
def call_ai_agent(provider, api_key, prompt):
    try:
        if provider == "DeepSeek (国内直连)":
            client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")
            response = client.chat.completions.create(
                model="deepseek-chat", messages=[{"role": "user", "content": prompt}], temperature=0.4
            )
            return response.choices[0].message.content
        else:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            response = model.generate_content(prompt)
            return response.text
    except Exception as e:
        return f"🔑 API 连接失败，请检查您的 Key 是否输入正确。错误信息：{str(e)}"

# ==========================================
# 4. 页面侧边栏布局
# ==========================================
with st.sidebar:
    st.markdown("### 🔑 1. 决策大脑接口配置")
    ai_provider = st.selectbox("大模型引擎选择", ["DeepSeek (国内直连)", "Google Gemini"])
    api_key = st.text_input(f"输入 {ai_provider} 密钥", type="password")
    
    st.markdown("### 📁 2. 混合研究文件导入区")
    uploaded_file = st.file_uploader("支持 CSV数据表 / PDF / Word / PPTX / TXT", type=["csv", "pdf", "docx", "pptx", "txt"])
    
    if uploaded_file and st.button("🚀 启动自动化分析流水线", use_container_width=True):
        file_ext = uploaded_file.name.split(".")[-1].lower()
        if file_ext == "csv":
            st.session_state.file_type = "dataset"
            df_cat, df_hier, r_dict = process_dataset_locally(uploaded_file)
            st.session_state.df_cleaned = df_cat
            st.session_state.df_hierarchical = df_hier
            st.session_state.diagnostic_report = r_dict
            st.session_state.analyzed = True
            # 清空历史，防止跨文件混淆
            st.session_state.chat_history_dataset = []
        else:
            st.session_state.file_type = "text"
            if not api_key:
                st.error("🔒 分析非结构化文本内容，必须先在上方配置您的 AI API Key！")
            else:
                with st.spinner("AI 正在对文本文件做结构化大纲提炼..."):
                    text_content = extract_text_from_doc(uploaded_file)
                    st.session_state.raw_text = text_content
                    prompt = f"请对以下学术/商业文本进行逻辑大纲提取（包含引言、研究对象、核心重点、数据结论或建议）：\n\n{text_content[:6000]}"
                    st.session_state.outline = call_ai_agent(ai_provider, api_key, prompt)
                    st.session_state.analyzed = True
                    st.session_state.chat_history_text = []
        st.rerun()

# ==========================================
# 5. 右侧主界面渲染模块
# ==========================================
st.markdown("<h1 style='text-align: center; color: #1E3A8A;'>🎓 SKU 商业清洗与 AI 科学研讨 Agent 看板</h1>", unsafe_allow_html=True)

if not st.session_state.analyzed:
    st.info("💡 请在左侧侧边栏导入您的原始数据集或文档，系统会自动识别并进入不同的专业看板。")
    st.image("https://img.icons8.com/illustrations/ul/480/searching.png", width=350)
else:
    # ------------------ 分流看板 A：数据集处理流 ------------------
    if st.session_state.file_type == "dataset":
        tab_data, tab_stat, tab_ai = st.tabs(["📊 A 本地数智清洗归档", "📐 B 统计检验前沿指引", "💬 C 穿透式 AI 答辩模拟"])
        
        with tab_data:
            st.markdown("### ⚡ 纯本地计算结果 (不消耗 API Token)")
            st.warning(st.session_state.diagnostic_report["warning"])
            
            # 本地图表绘制
            col_g1, col_g2 = st.columns(2)
            df_c = st.session_state.df_cleaned
            with col_g1:
                fig1, ax1 = plt.subplots(figsize=(6, 4))
                ax1.pie(df_c['total_revenue'], labels=df_c['category'], autopct='%1.1f%%', startangle=90)
                ax1.set_title("品类宏观整体利润占比")
                st.pyplot(fig1)
                plt.close(fig1)
            with col_g2:
                fig2, ax2 = plt.subplots(figsize=(6, 4))
                comp_df = pd.concat([st.session_state.diagnostic_report["top_cat"], st.session_state.diagnostic_report["bottom_cat"]])
                ax2.barh(comp_df['category'], comp_df['total_revenue'], color=['green']*3 + ['red']*3)
                ax2.set_title("宏观优劣品类两极分化测度")
                ax2.invert_yaxis()
                st.pyplot(fig2)
                plt.close(fig2)
                
            st.markdown("#### 📂 穿透级联报表下载区 (已自动创建并归档)")
            col_d1, col_d2 = st.columns(2)
            col_d1.download_button(
                label="📥 导出格式化：清洗后文件.csv", data=st.session_state.df_cleaned.to_csv(index=False).encode('utf-8-sig'),
                file_name="清洗后文件.csv", mime="text/csv"
            )
            col_d2.download_button(
                label="📥 导出层级穿透大表：分析后文件.csv", data=st.session_state.df_hierarchical.to_csv(index=False).encode('utf-8-sig'),
                file_name="分析后文件.csv", mime="text/csv"
            )
            st.dataframe(st.session_state.df_hierarchical, height=350, use_container_width=True)

        with tab_stat:
            st.markdown("### 📐 数据集定量特征与建模指向")
            st.info(st.session_state.diagnostic_report["stat_advise"])
            st.markdown("""
            **📢 研究Feasibility自查指引：**
            * **独立样本 T 检验**：适用于检验分类字段在仅有两组别时（如：线上 vs 线下），连续变量均值是否存在显著异同。
            * **方差分析 (ANOVA)**：用于多组别（如：A类、B类、C类）对销售额影响的显著性差异探测。
            * **OLS 回归**：若数字指标多维共存，用于测度自变量对因变量贡献的因果关联。
            """)

        with tab_ai:
            st.markdown("### 💬 穿透式数据决策 AI 互动")
            if not api_key:
                st.error("🔒 请在侧边栏输入您的 API Key 以激活 AI 对话功能。")
            else:
                # 渲染数据集独有的对话流
                for msg in st.session_state.chat_history_dataset:
                    with st.chat_message(msg["role"]): 
                        st.markdown(msg["content"])
                
                # 为该输入框设定全局唯一 Key：input_dataset
                dataset_input = st.chat_input("针对刚才生成的清洗报告或统计检验方向，对 AI 提问...", key="input_dataset")
                if dataset_input:
                    with st.chat_message("user"): 
                        st.markdown(dataset_input)
                    st.session_state.chat_history_dataset.append({"role": "user", "content": dataset_input})
                    
                    with st.chat_message("assistant"):
                        with st.spinner("AI 正在深度透视底层数据并生成商业故事线..."):
                            context = f"""
                            你是一个极度高级的商业咨询与学术答辩专家。
                            刚刚我们通过 Python 本地算法完成了数据集分析：
                            - 诊断结论：{st.session_state.diagnostic_report['warning']}
                            - 推荐的检验方向：{st.session_state.diagnostic_report['stat_advise']}
                            - 利润最高的单品ID为: {st.session_state.diagnostic_report['top_prod']['product_id'].tolist()}
                            
                            请基于精益创业和科学实证的逻辑，针对用户的提问进行解答，提供答辩话术或优化建策：
                            用户问题：{dataset_input}
                            """
                            reply = call_ai_agent(ai_provider, api_key, context)
                            st.markdown(reply)
                    st.session_state.chat_history_dataset.append({"role": "assistant", "content": reply})
                    st.rerun()

    # ------------------ 分流看板 B：传统文本处理流 ------------------
    else:
        tab_text, tab_chat = st.tabs(["📋 A 结构化文献大纲", "💬 B 文本理论深度研讨"])
        with tab_text:
            st.markdown("### 📑 AI 提炼的核心架构")
            st.markdown(st.session_state.outline)
        with tab_chat:
            st.markdown("### 💬 学术论文/课件出题与辩论")
            # 渲染文本类独有的对话流
            for msg in st.session_state.chat_history_text:
                with st.chat_message(msg["role"]): 
                    st.markdown(msg["content"])
                    
            # 为该输入框设定全局唯一 Key：input_text
            text_input = st.chat_input("请输入您想针对本篇课件/论文探讨的概念或模拟出题...", key="input_text")
            if text_input:
                with st.chat_message("user"): 
                    st.markdown(text_input)
                st.session_state.chat_history_text.append({"role": "user", "content": text_input})
                
                with st.chat_message("assistant"):
                    with st.spinner("AI 正在检索上下文..."):
                        context = f"你是一个学术助教。基于用户上传的文本片段：{st.session_state.raw_text[:3000]}。请回答：{text_input}"
                        reply = call_ai_agent(ai_provider, api_key, context)
                        st.markdown(reply)
                st.session_state.chat_history_text.append({"role": "assistant", "content": reply})
                st.rerun()
